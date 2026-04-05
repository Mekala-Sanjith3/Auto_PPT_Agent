"""
MCP server: PowerPoint creation with mandatory outline planning before slides.
Produces styled decks with themed backgrounds, accent bars, and image+bullet layouts.
The server enforces a plan-first workflow: submit_slide_plan() MUST be called before
any add_slide() or add_slide_with_image() call - this is the "agentic loop" design.
"""

#allow newer type hint syntax (e.g. list[dict]) on older Python versions
from __future__ import annotations

import json     #parse bullet arrays and slide plan JSON sent by the LLM
import re       # sanitize filenames and normalize title strings for comparison
import sys      #modify sys.path so we can import from the project root folder
from pathlib import Path         #cross-platform file path handling
from typing import Any, Optional, Tuple  #type hints for older Python compatibility

# PIL: used to read image dimensions before placing images (for contain-fit math)
from PIL import Image as PILImage

#add project root to Python path so we can import config.py from here
# needed because this file lives in servers/ but config.py is one level up.
ROOT = Path(__file__).resolve().parent.parent   # Navigate up:servers/ → Agent_PPT/
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))   #Ensure project root is first in path

# FastMCP: the framework that makes Python functions into MCP-callable tools
from mcp.server.fastmcp import FastMCP

# python-pptx: the library that actually creates and saves .pptx files
from pptx import Presentation
from pptx.dml.color import RGBColor      # represents colors as RGB hex values
from pptx.enum.shapes import MSO_SHAPE  # Enum of built-in shape types (rectangle, oval, etc.)
from pptx.util import Inches, Pt, Emu   #unit converters: Inches(), Pt() for font size, Emu() for internal units

# Shared config: output directory, HF token, etc.
from config import Config


#themes dictionary: each theme defines 8 colors used throughout the deck.
# Keys: bg (background), bg_alt (alternate background), text (main text),
#       muted (secondary/label text), accent (primary accent bar color),
#       accent2 (secondary decorative color), frame (image border), num (slide number)
THEMES = {
    # napkin: warm cream + coral - friendly, modern, hand-drawn feel (the default theme)
    "napkin": {
        "bg":      RGBColor(0xFC, 0xFA, 0xF7),  #off-white cream -warm and easy on the eyes
        "bg_alt":  RGBColor(0xF5, 0xF3, 0xEF),  # Slightly darker cream - used for alternate areas
        "text":    RGBColor(0x2D, 0x34, 0x36),  # Near-black charcoal- high contrast on cream
        "muted":   RGBColor(0x63, 0x6E, 0x72),  # Medium grey - used for labels and secondary text
        "accent":  RGBColor(0xFF, 0x6B, 0x6B),  # Coral-red - left accent bar and decorative dots
        "accent2": RGBColor(0x00, 0xCE, 0xAC),  # Teal - secondary decorative elements
        "frame":   RGBColor(0xE8, 0xE4, 0xDE),  # Light warm grey - image border frame
        "num":     RGBColor(0xAA, 0xA5, 0x9F),  # Soft taupe - slide number text
    },
    #ocean: deep navy + teal - professional and bold, good for technical topics
    "ocean": {
        "bg":      RGBColor(0x0D, 0x21, 0x37),  # Very dark navy - dramatic dark background
        "bg_alt":  RGBColor(0x15, 0x2B, 0x45),  # Slightly lighter navy - for alternate sections
        "text":    RGBColor(0xCC, 0xEE, 0xFF),  # Light sky blue - readable on dark background
        "muted":   RGBColor(0x64, 0xDF, 0xDF),  # Cyan-teal - secondary labels on dark bg
        "accent":  RGBColor(0x00, 0xCE, 0xAC),  # Teal - accent bar
        "accent2": RGBColor(0xFF, 0xA5, 0x00),  # Amber orange - contrasting highlight
        "frame":   RGBColor(0x20, 0x3A, 0x55),  # Dark steel blue - subtle image border
        "num":     RGBColor(0x40, 0x70, 0x90),  # Muted blue - slide numbers
    },
    # dark: deep indigo + electric red - dramatic, high-impact look
    "dark": {
        "bg":      RGBColor(0x1A, 0x1A, 0x2E),  # Deep indigo-black - rich dark background
        "bg_alt":  RGBColor(0x22, 0x22, 0x38),  # Slightly lighter indigo - alternate areas
        "text":    RGBColor(0xE0, 0xE0, 0xFF),  # Lavender-white - softer than pure white on dark bg
        "muted":   RGBColor(0x88, 0x88, 0xAA),  # Muted purple-grey - secondary text
        "accent":  RGBColor(0xE9, 0x45, 0x60),  # Electric red-pink - bold accent bar
        "accent2": RGBColor(0x4A, 0x9F, 0xFF),  # Bright blue - secondary accent
        "frame":   RGBColor(0x30, 0x30, 0x50),  # Dark purple - subtle image border
        "num":     RGBColor(0x55, 0x55, 0x88),  # Purple-grey - slide numbers
    },
    #minimal: clean white + blue - corporate and professional
    "minimal": {
        "bg":      RGBColor(0xFF, 0xFF, 0xFF),  # Pure white - clean, no-frills background
        "bg_alt":  RGBColor(0xF8, 0xF8, 0xF8),  # Near-white - very subtle alternate background
        "text":    RGBColor(0x1A, 0x1A, 0x1A),  # Near-black - maximum contrast on white
        "muted":   RGBColor(0x66, 0x66, 0x66),  # Medium grey — secondary text
        "accent":  RGBColor(0x00, 0x7A, 0xFF),  # Bright blue - clean corporate accent
        "accent2": RGBColor(0xFF, 0x6B, 0x35),  # Orange - secondary highlight
        "frame":   RGBColor(0xDD, 0xDD, 0xDD),  # Light grey - subtle image border
        "num":     RGBColor(0xBB, 0xBB, 0xBB),  # Light grey - slide numbers
    },
}


# pptx_instructions: shown to Claude as the tool usage guide.
# This enforces the plan-first agentic loop pattern required by the assignment.
PPTX_INSTRUCTIONS = """
You MUST follow this exact order:
(1) create_presentation(topic, filename_base) - optionally pass theme_name
(2) submit_slide_plan(JSON array of {title, bullets}) - plan ALL slides FIRST before any add_slide
(3) For each content slide (recommended):
    a) search_topic(query) via auto-ppt-web - get real facts
    b) get_image_for_slide(prompt) via auto-ppt-hf-image - get illustration path
    c) add_slide_with_image(title, bullets_json, image_path)
   OR use add_slide(title, bullets_json) if no image is needed.
(4) save_presentation()

Rules:
- NEVER call add_slide before submit_slide_plan.
- Aim for 3-5 bullets per slide, concise and factual.
- Use add_slide_with_image for most slides - produces napkin-style layout.
- Available themes: napkin (default), ocean, dark, minimal.
"""

#create the FastMCP server instance with the instruction string above
mcp = FastMCP("AutoPPT", instructions=PPTX_INSTRUCTIONS)


#session state: module-level variables that hold the current presentation.
# These are reset each time create_presentation() is called.
# Using module globals is intentional — this is a single-session MCP server.
_prs:             Optional[Presentation] = None   # the python-pptx Presentation object in memory
_output_path:     Optional[Path]         = None   #full path where the .pptx will be saved
_topic:           str                    = ""     # the presentation topic (shown on title slide + corner tag)
_slide_plan:      Optional[list[dict]]   = None   #the approved outline (list of {title, bullets} dicts)
_slides_added:    int                    = 0      # counter: how many content slides have been added so far
_plan_submitted:  bool                   = False  # gate flag: True after submit_slide_plan succeeds
_theme:           dict                   = THEMES["napkin"]  # Active color dictionary for the current deck
_theme_name:      str                    = "napkin"          # name of the active theme (for status/save messages)
_total_planned:   int                    = 0      # total number of slides in the approved plan


def _reset_session() -> None:
    """
    resets all session-state globals back to their initial values.
    Called at the start of create_presentation() to ensure a clean slate
    when starting a new presentation (prevents bleed-over from previous calls).
    """
    global _prs, _output_path, _topic, _slide_plan, _slides_added
    global _plan_submitted, _theme, _theme_name, _total_planned

    _prs            = None             #Clear the in-memory presentation object
    _output_path    = None             # Clear the target save path
    _topic          = ""               # Clear the topic string
    _slide_plan     = None             #Clear the approved outline
    _slides_added   = 0                #Reset the slide counter
    _plan_submitted = False            # Reset the planning gate
    _theme          = THEMES["napkin"] #Reset to default theme colors
    _theme_name     = "napkin"         # Reset to default theme name
    _total_planned  = 0                # Reset the planned slide count


def _norm(s: str) -> str:
    """
    Normalizes a string for loose comparison: lowercase, collapse whitespace.
    used to compare the title the LLM passed in vs the planned title,
    so minor capitalization or spacing differences don't cause false mismatches.
    """
    return re.sub(r"\s+", " ", s.strip().lower())   # Collapse all whitespace runs to single space


def _strip_placeholder_shapes(slide) -> None:
    """
    Removes all default placeholder shapes (title/content boxes) from a slide.
    called before building slides from scratch so we have a fully blank canvas.
    the blank layout (index 6) sometimes still includes hidden placeholders - this clears them.
    errors are caught per-shape so one bad shape doesn't prevent the rest from being removed.
    """
    for shape in list(slide.shapes):   # list() creates a copy so we can safely remove during iteration
        try:
            if getattr(shape, "is_placeholder", False):  # Only remove placeholder shapes, not regular shapes
                el = shape.element            # Get the underlying XML element
                el.getparent().remove(el)     # Remove it from the slide's XML tree
        except Exception:
            continue   # skip this shape if something went wrong - don't abort the whole operation


def _contain_picture_layout(
    image_path: str,
    box_left_in: float, box_top_in: float,
    box_w_in: float,    box_h_in: float,
) -> Tuple[float, float, float]:
    """
    Calculates the position and width to place an image using contain-fit scaling.
    'Contain-fit' means the image is scaled to fit entirely within the box
    while maintaining its original aspect ratio - no cropping, no stretching.
    Centers the image within the box.

    Returns: (left_in, top_in, width_in) - all in inches for python-pptx.
    Falls back to the box dimensions if the image cannot be read.
    """
    try:
        with PILImage.open(image_path) as im:
            iw, ih = im.size   # Read actual pixel dimensions of the image
    except Exception:
        return box_left_in, box_top_in, box_w_in   # can't read image - return box as-is

    if iw <= 0 or ih <= 0:
        return box_left_in, box_top_in, box_w_in   # degenerate image dimensions - bail out

    ar = iw / ih         # image aspect ratio (width / height)
    bw, bh = box_w_in, box_h_in   # Box width and height in inches
    box_ar = bw / bh     # Box aspect ratio

    # determine whether to fit by width or by height
    if ar > box_ar:
        # image is wider than the box relative to height - fit width, letterbox top/bottom
        w_in = bw
        h_in = bw / ar   # Scale height proportionally
    else:
        # image is taller than the box relative to width - fit height, pillarbox left/right
        h_in = bh
        w_in = bh * ar   # Scale width proportionally

    # center the scaled image within the box by offsetting from the box's top-left corner
    left_in = box_left_in + (bw - w_in) / 2   # Horizontal centering offset
    top_in  = box_top_in  + (bh - h_in) / 2   # Vertical centering offset

    return left_in, top_in, w_in   # return position and width


def _apply_background(slide) -> None:
    """
    sets the slide's background to the current theme's background color.
    We first try to detach from the master slide's background so our custom
    color isn't overridden. then we set a solid fill with the theme bg color.
    Both operations are wrapped in try/except because background access
    can fail on some slide layout types.
    """
    try:
        slide.follow_master_background = False   # stop inheriting background from the slide master
    except Exception:
        pass   # not all slide types support this property - ignore if it fails

    try:
        fill = slide.background.fill   # Access the background fill object
        fill.solid()                   # Switch to solid color fill mode
        fill.fore_color.rgb = _theme["bg"]   # Apply the theme background color
    except Exception:
        pass   # if the background API fails (some layouts), leave it as-is


def _add_accent_bar(slide) -> None:
    """
    Adds a narrow vertical rectangle on the left edge of the slide.
    this is the signature visual element of the napkin-style layout -
    a bold colored stripe that gives the slide a branded look.
    """
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,   # Shape type: plain rectangle
        Inches(0),             # Left edge: flush with the left side of the slide
        Inches(0),             # Top edge: start at the very top
        Inches(0.10),          # Width: 0.1 inch - narrow stripe
        Inches(7.5),           # Height: full slide height (standard is 7.5")
    )
    bar.fill.solid()                       # Solid color fill
    bar.fill.fore_color.rgb = _theme["accent"]  # use the theme accent color
    bar.line.fill.background()             # Remove the shape's border (transparent line)


def _add_slide_number(slide, number: int, total: int) -> None:
    """
    adds a small slide-number label in the bottom-right corner: e.g. '2 / 5'.
    The number includes the title slide in the count (hence +1 in the caller).
    """
    label = f"{number} / {total}"   # Format: current slide / total slides

    # Add a small text box in the bottom-right corner of the slide
    tb = slide.shapes.add_textbox(
        Inches(9.1),   # Left: near the right edge (slide is 10" wide)
        Inches(7.05),  # Top: near the bottom edge (slide is 7.5" tall)
        Inches(0.8),   # Width: enough for "12 / 20"
        Inches(0.3),   # Height: single line
    )
    tf = tb.text_frame
    tf.text = label            # Set the number text

    p = tf.paragraphs[0]       # Get the first (only) paragraph
    p.font.size      = Pt(9)   # Small font - unobtrusive
    p.font.color.rgb = _theme["num"]   # use the theme's number color (usually muted)
    p.font.name      = "Calibri"       # Match the rest of the slide typography


def _add_divider_line(slide, top_in: float) -> None:
    """
    Draws a thin horizontal rule below the slide title to visually separate
    the title area from the content area. Uses a very flat rectangle as a line.
    top_in: vertical position in inches from the top of the slide.
    """
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.50),     # Start slightly inset from the left (after the accent bar)
        Inches(top_in),   # Vertical position passed by the caller
        Inches(9.0),      # Width: spans most of the slide horizontally
        Inches(0.025),    # Height: very thin - acts as a visual line
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _theme["frame"]  # subtle divider color from theme
    line.line.fill.background()                 # No border on the divider shape itself


def _add_topic_tag(slide, topic: str) -> None:
    """
    Adds a small italic topic label in the top-right corner of the slide.
    Truncated to 28 characters with an ellipsis if longer.
    Gives viewers context about the overall presentation topic on every slide.
    """
    short = topic[:28] + ("…" if len(topic) > 28 else "")   # Truncate long topics with ellipsis

    # Small text box in the top-right corner
    tb = slide.shapes.add_textbox(
        Inches(6.8),   # Start 6.8" from left - leaves room for the slide title
        Inches(0.12),  # Near the top
        Inches(3.0),   # Width: 3 inches
        Inches(0.3),   # Height: single line
    )
    tf = tb.text_frame
    tf.text = short

    p = tf.paragraphs[0]
    p.font.size      = Pt(8)           # Very small - purely informational
    p.font.color.rgb = _theme["muted"] # Muted color so it doesn't compete with the title
    p.font.name      = "Calibri"
    p.font.italic    = True            # Italic distinguishes it from slide content


def _style_title_slide(slide, topic: str, subtitle: str = "Learning deck · visual outline") -> None:
    """
    applies styling to the first slide (title slide):
    - Background color from theme
    - Large bold title text
    - Subtitle text in muted color
    - Three decorative dots in the right half
    - Accent bar on the left edge
    The title and subtitle placeholders come from the slide layout (index 0).
    """
    _apply_background(slide)           # Fill the background with the theme color
    slide.shapes.title.text = topic    # Set the main title text from the topic argument

    # Style the title placeholder: large, bold, themed color
    tf = slide.shapes.title.text_frame
    tf.margin_left = Inches(0.08)      # Small left margin so text isn't flush with the accent bar
    for p in tf.paragraphs:
        p.font.size      = Pt(40)      # Large title font - 40pt is presentation-ready
        p.font.bold      = True
        p.font.color.rgb = _theme["text"]  # Main text color from theme
        p.font.name      = "Calibri"

    # Style the subtitle placeholder (placeholder index 1) if it exists in this layout
    if len(slide.placeholders) > 1:
        try:
            st = slide.placeholders[1]   # The subtitle placeholder
            st.text = subtitle           # Set the subtitle text
            for p in st.text_frame.paragraphs:
                p.font.size      = Pt(18)          # Smaller than the title
                p.font.color.rgb = _theme["muted"] # Muted/secondary color
                p.font.name      = "Calibri"
        except Exception:
            pass   # some layouts don't have a subtitle placeholder - silently skip


    # decorative dots: three overlapping circles in the right portion of the slide.
    # these add visual interest and a 'napkin doodle' feel to the title slide.
        # each tuple is (left_position, top_position, theme_color_key).
    dots = [
        (Inches(8.1), Inches(3.8), "accent2"),   # Top-right dot: teal
        (Inches(8.8), Inches(4.4), "accent"),    # Middle-right dot: coral
        (Inches(7.8), Inches(4.9), "accent2"),   # Lower-right dot: teal
    ]
    for left, top, key in dots:
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,          # Circle shape (oval with equal width/height)
            left, top,
            Inches(0.38),            # Width: small circle
            Inches(0.38),            # Height: same as width - perfect circle
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = _theme[key]   # use the color from the THEMES dict
        dot.line.fill.background()              # No border on the decorative dots

    # Add the signature left accent bar to the title slide as well
    _add_accent_bar(slide)


# MCP tools - these are the functions Claude calls during the agentic loop.
# each is decorated with @mcp.tool() to register it as a callable MCP tool.

@mcp.tool()
def create_presentation(
    topic: str,
    filename_base: str,
    theme_name: str = "napkin",
) -> str:
    """
    Initializes a new PowerPoint presentation in memory and adds the title slide.
    must be called first - before submit_slide_plan or any add_slide call.

    topic: the main subject (shown on the title slide and in the corner of every slide)
    filename_base: the filename without extension (special chars are sanitized automatically)
    theme_name: visual style - one of: napkin (default), ocean, dark, minimal
    """
    # Declare all session globals we'll be modifying
    global _prs, _output_path, _topic, _slide_plan, _slides_added
    global _plan_submitted, _theme, _theme_name, _total_planned

    _reset_session()   # Clear any data from a previous create_presentation call

    # Validate and apply the requested theme (fall back to napkin if unknown)
    chosen = theme_name.lower().strip()
    if chosen not in THEMES:
        chosen = "napkin"    # unknown theme name - silently fall back to the default
    _theme      = THEMES[chosen]   # Load the color dictionary for the chosen theme
    _theme_name = chosen           # Store the theme name for status messages

    # Ensure the output directory exists before computing the save path
    Config.ensure_output_dir()

    # sanitize the filename: keep only alphanumeric, hyphen, and underscore characters
    # truncate to 80 chars to avoid filesystem path-length issues on Windows
    safe = re.sub(r"[^\w\-]", "_", filename_base)[:80] or "presentation"
    _output_path = Config.OUTPUT_DIR / f"{safe}.pptx"   # Full absolute path for saving later
    _topic = topic   # Store the topic for use in slide corner tags and status messages

    # create the Presentation object (starts with the default blank template)
    _prs = Presentation()

    # add the title slide using slide layout index 0 (Title Slide layout from the template)
    title_layout = _prs.slide_layouts[0]
    slide = _prs.slides.add_slide(title_layout)
    _style_title_slide(slide, topic)   # Apply background, fonts, decorative elements

    return (
        f"✓ Created presentation '{topic}' (theme: {chosen}). "
        f"Saves to: {_output_path}. "
        "REQUIRED next step: call submit_slide_plan with your full JSON outline."
    )


@mcp.tool()
def submit_slide_plan(plan_json: str) -> str:
    """
    Validates and stores the slide outline before any content slides are added.
    This is the mandatory planning gate - no add_slide call is allowed until this succeeds.

    plan_json: a JSON array of objects, each with:
        - "title": string (required)
        - "bullets": array of 3-5 strings (required)

    example: [{"title": "What is a Star?", "bullets": ["Nuclear fusion...", "..."]}, ...]
    """
    global _slide_plan, _plan_submitted, _slides_added, _total_planned

    # Guard: make sure create_presentation was called first
    if _prs is None:
        return "Error: call create_presentation first."

    # Parse the JSON string into a Python object
    try:
        data = json.loads(plan_json)
    except json.JSONDecodeError as e:
        return f"Invalid JSON: {e}. Fix and retry."   # tell the LLM exactly what went wrong

    # validate top-level structure: must be a non-empty list
    if not isinstance(data, list) or len(data) == 0:
        return "plan_json must be a non-empty JSON array of slide objects."

    # Validate and clean each slide object
    cleaned: list[dict] = []
    for i, item in enumerate(data):
        if not isinstance(item, dict):
            return f"Item {i} must be an object with 'title' and 'bullets'."

        title   = str(item.get("title", "")).strip()   # Extract and strip the title string
        bullets = item.get("bullets", [])              # Extract the bullets field

        # normalize bullets: handle if the LLM mistakenly sent a string instead of a list
        if isinstance(bullets, str):
            bullets = [bullets]

        # Convert all items to stripped strings and remove empty ones
        bullets = [
            str(b).strip()
            for b in (bullets if isinstance(bullets, list) else [])
            if str(b).strip()
        ]

        # Validate: title is required
        if not title:
            return f"Item {i} needs a non-empty title."

        # if no valid bullets were provided, add a placeholder so the slide isn't empty
        if not bullets:
            bullets = ["(Placeholder - fill from topic knowledge.)"]

        cleaned.append({
            "title":   title,
            "bullets": bullets[:8],   # Cap at 8 bullets per slide to avoid overcrowding
        })

    # Store the validated plan and set the gate flag to allow add_slide calls
    _slide_plan     = cleaned
    _plan_submitted = True       # unlock the add_slide and add_slide_with_image tools
    _slides_added   = 0          # reset counter (in case of re-planning)
    _total_planned  = len(cleaned)   # store total for slide-number labels and completion checks

    # build a preview of the slide order to confirm to the LLM
    titles = [c["title"] for c in cleaned]
    return (
        f"✓ Outline accepted: {_total_planned} slides planned.\n"
        "slide order: " + " - ".join(titles) + "\n"
        "Now call add_slide or add_slide_with_image for each slide, then save_presentation."
    )


@mcp.tool()
def add_slide(title: str, bullets_json: str) -> str:
    """
    Adds a text-only content slide (no image).
    Use this when get_image_for_slide is not available or not needed.
    Applies background, accent bar, styled title, divider, and formatted bullets.

    title: slide heading (should match the planned title from submit_slide_plan)
    bullets_json: JSON array of 3-5 bullet point strings
    must call submit_slide_plan first - will error otherwise.
    """
    global _slides_added

    # Guard: presentation must be initialized
    if _prs is None:
        return "Error: call create_presentation first."

    # guard: the plan must be submitted before any slides can be added (the agentic gate)
    if not _plan_submitted or not _slide_plan:
        return "Error: call submit_slide_plan with your full outline before add_slide."

    # guard: don't add more slides than were planned
    if _slides_added >= len(_slide_plan):
        return "All planned slides already added. Call save_presentation."

    # Parse the bullets JSON (handle edge cases where the LLM sends a string)
    try:
        bullets = json.loads(bullets_json)
        if not isinstance(bullets, list):
            bullets = [str(bullets)]    # Wrap non-list in a list
        bullets = [str(b).strip() for b in bullets if str(b).strip()]  # Clean empty strings
    except json.JSONDecodeError:
        bullets = [bullets_json]   # If JSON parsing fails entirely, treat the whole string as one bullet
    if not bullets:
        bullets = ["(No bullet content — expand from topic.)"]   # Placeholder if nothing valid

    # Check if the title matches what was planned (loose comparison, case-insensitive)
    expected = _slide_plan[_slides_added]["title"]
    note = "" if _norm(title) == _norm(expected) else \
        f"Note: title differs from planned '{expected}'; using '{title}'. "

    # Add a slide using layout index 1 (Title and Content layout)
    layout = _prs.slide_layouts[1]
    slide  = _prs.slides.add_slide(layout)

    _apply_background(slide)   # Apply theme background color
    _add_accent_bar(slide)     # Draw the left vertical accent stripe

    # Set and style the title text
    slide.shapes.title.text = title
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.size      = Pt(28)          # Slightly smaller than title slide (28pt vs 40pt)
        p.font.bold      = True
        p.font.color.rgb = _theme["text"]
        p.font.name      = "Calibri"

    # Draw the thin horizontal rule below the title
    _add_divider_line(slide, 1.10)

    # add bullet points to the body placeholder
    body = slide.placeholders[1]   # index 1 is always the content/body placeholder
    tf   = body.text_frame
    tf.clear()           # remove any default text from the layout
    tf.word_wrap = True  # enable word wrapping so long bullets don't overflow

    for i, b in enumerate(bullets[:8]):   # limit to 8 bullets maximum
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text           = b
        p.level          = 0
        p.font.size      = Pt(18)          # Larger, more readable body text
        p.font.color.rgb = _theme["text"]
        p.font.name      = "Calibri"
        p.space_before   = Pt(6)
        p.space_after    = Pt(6)
        p.line_spacing   = 1.3

    # Slide number only — NO topic tag
    _slides_added += 1
    _add_slide_number(slide, _slides_added + 1, _total_planned + 1)

    return note + f"✓ Added text slide {_slides_added}/{_total_planned}: '{title}'"


@mcp.tool()
def add_slide_with_image(
    title: str,
    bullets_json: str,
    image_path: str,
) -> str:
    """
    Adds a styled content slide with text on the left and an image on the right.
    This is the preferred tool - produces the full napkin-style layout with:
    - Background + accent bar
    - Title with divider
    - Bullet column on the left (4.4" wide)
    - Image column on the right with a rounded frame
    - Topic tag and slide number

    title: slide heading
    bullets_json: JSON array of bullet strings
    image_path: absolute file path to a PNG or JPG (returned by get_image_for_slide)
    Must call submit_slide_plan first.
    """
    global _slides_added

    # guard: presentation must exist
    if _prs is None:
        return "Error: call create_presentation first."

    #guard: planning gate must be passed
    if not _plan_submitted or not _slide_plan:
        return "Error: call submit_slide_plan first."

    # Guard: don't exceed the planned slide count
    if _slides_added >= len(_slide_plan):
        return "All planned slides already added."

    # parse the bullets (same logic as add_slide)
    try:
        bullets = json.loads(bullets_json)
        if not isinstance(bullets, list):
            bullets = [str(bullets)]
        bullets = [str(b).strip() for b in bullets if str(b).strip()]
    except json.JSONDecodeError:
        bullets = [bullets_json]
    if not bullets:
        bullets = ["(Visual emphasis — see image.)"]   # Default when no bullets provided

    # Resolve the image path and verify the file actually exists
    path = Path(image_path).expanduser().resolve()
    if not path.is_file():
        return f"Image not found at {path}. Use add_slide without image instead."

    # use the blank layout (index 6) so we can place all shapes manually with full control
    #fall back to layout 1 if the template has fewer than 7 layouts
    blank_layout = _prs.slide_layouts[6] if len(_prs.slide_layouts) > 6 else _prs.slide_layouts[1]
    slide = _prs.slides.add_slide(blank_layout)

    _apply_background(slide)           # Apply theme background color
    _add_accent_bar(slide)             # Draw the left accent stripe
    _strip_placeholder_shapes(slide)   # Remove any hidden placeholder shapes from the layout

    # ── Layout constants  (all in inches) ─────────────────────────────────────
    # Two-column grid: title spans full width at top;
    # left 48% = bullet text, right 47% = image filling the full content height.

    MARGIN_L    = 0.45
    TITLE_TOP   = 0.18   # title sits close to the top — more content room
    TITLE_H     = 0.80   # tall enough for wrapping titles
    DIVIDER_TOP = 1.03   # thin rule just below the title box
    BODY_TOP    = 1.18   # bullets start right after the rule
    BODY_H      = 5.82   # bullets fill nearly all remaining vertical space
    TEXT_COL_W  = 4.55   # left text column (~48% of 10" slide)
    COL_GAP     = 0.20   # small gap between text and image columns
    IMG_LEFT    = MARGIN_L + TEXT_COL_W + COL_GAP  # 0.45+4.55+0.20 = 5.20"
    IMG_W       = 10.0 - IMG_LEFT - 0.25           # 10.0-5.20-0.25 = 4.55" wide
    IMG_TOP     = BODY_TOP                          # image top aligns with first bullet
    IMG_H       = BODY_H                            # image fills the full content height
    PAD         = 0.10   # breathing room inside the image frame

    # ── Place the image with contain-fit scaling ────────────────────────────
    pl, pt, pw = _contain_picture_layout(
        str(path),
        IMG_LEFT + PAD,
        IMG_TOP  + PAD,
        IMG_W    - 2 * PAD,
        IMG_H    - 2 * PAD,
    )
    try:
        slide.shapes.add_picture(str(path), Inches(pl), Inches(pt), width=Inches(pw))
    except Exception as e:
        return f"Image insert failed: {e}. Use add_slide without image."

    # Subtle rounded-rect frame around the image column area
    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(IMG_LEFT), Inches(IMG_TOP),
        Inches(IMG_W),    Inches(IMG_H),
    )
    frame.fill.background()
    frame.line.color.rgb = _theme["frame"]
    frame.line.width     = Inches(0.010)  # very thin border

    # Move frame behind the image in the XML tree
    sp = frame._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

    # ── Title text box (full-width) ─────────────────────────────────────────
    title_box = slide.shapes.add_textbox(
        Inches(MARGIN_L),
        Inches(TITLE_TOP),
        Inches(10.0 - MARGIN_L - 0.25),  # near-full width
        Inches(TITLE_H),
    )
    ttf = title_box.text_frame
    ttf.word_wrap = True
    ttf.text = title
    for p in ttf.paragraphs:
        p.font.size      = Pt(28)
        p.font.bold      = True
        p.font.color.rgb = _theme["text"]
        p.font.name      = "Calibri"

    # Thin divider line below title
    _add_divider_line(slide, DIVIDER_TOP)

    # ── Bullet points text box (left column) ────────────────────────────────
    body_box = slide.shapes.add_textbox(
        Inches(MARGIN_L),
        Inches(BODY_TOP),
        Inches(TEXT_COL_W),
        Inches(BODY_H),
    )
    tf = body_box.text_frame
    tf.word_wrap = True

    BULLET_CHAR = "\u25cf"  # ● solid circle — cleaner than "•"

    for i, b in enumerate(bullets[:8]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Accent-coloured bullet glyph in a Run, then text in a second Run
        run_dot = p.add_run()
        run_dot.text           = BULLET_CHAR + " "
        run_dot.font.size      = Pt(11)
        run_dot.font.color.rgb = _theme["accent"]
        run_dot.font.name      = "Calibri"
        run_dot.font.bold      = False

        run_txt = p.add_run()
        run_txt.text           = str(b)
        run_txt.font.size      = Pt(15)
        run_txt.font.color.rgb = _theme["text"]
        run_txt.font.name      = "Calibri"
        run_txt.font.bold      = False

        p.space_before = Pt(4)
        p.space_after  = Pt(10)
        p.line_spacing = 1.3

    # ── Slide number only — NO topic tag ────────────────────────────────────
    _slides_added += 1
    _add_slide_number(slide, _slides_added + 1, _total_planned + 1)

    return f"✓ Added image slide {_slides_added}/{_total_planned}: '{title}'"


@mcp.tool()
def save_presentation() -> str:
    """
    Saves the in-memory Presentation object to a .pptx file on disk.
    Call this after all add_slide / add_slide_with_image calls are complete.
    will warn (but still save) if fewer slides were added than planned.
    """
    # guard: nothing to save if create_presentation was never called
    if _prs is None or _output_path is None:
        return "Error: nothing to save - call create_presentation first."

    # Generate a warning if the LLM added fewer slides than it planned
    warn = ""
    if _slide_plan and _slides_added < len(_slide_plan):
        warn = (
            f"Warning: only {_slides_added} of {len(_slide_plan)} planned slides added. "
            "Consider completing remaining slides before saving. Saving anyway. "
        )

    try:
        Config.ensure_output_dir()   # make sure the output folder still exists (safety check)
        _prs.save(str(_output_path)) # Write the .pptx file to disk using the stored path
        return (
            warn
            + f"✓ Saved: {_output_path} "
            + f"({_slides_added + 1} slides total, theme: {_theme_name})"
            # +1 in the count because the title slide is not tracked in _slides_added
        )
    except Exception as e:
        return f"Save failed: {e}. Check AUTO_PPT_OUTPUT_DIR permissions."


@mcp.tool()
def get_status() -> str:
    """
    Returns the current session state as a single-line debug string.
    Useful for Claude (or a developer) to check progress mid-session.
    shows: topic, theme, whether the plan was submitted, how many slides were added vs planned, and the output path.
    """
    planned = len(_slide_plan) if _slide_plan else 0  # avoid None length error
    return (
        f"topic='{_topic}' | theme={_theme_name} | "
        f"plan_submitted={_plan_submitted} | "
        f"slides_added={_slides_added}/{planned} | "
        f"output={_output_path}"
    )


@mcp.tool()
def list_themes() -> str:
    """
    Returns a formatted list of all available theme names with descriptions.
    The LLM can call this before create_presentation to choose the right theme
    for the topic (e.g. 'ocean' for science, 'napkin' for general education).
    """
    descriptions = {
        "napkin":  "warm cream canvas, coral accent — friendly & modern (default)",
        "ocean":   "deep navy background, teal accent — professional & bold",
        "dark":    "deep indigo background, electric red accent — dramatic",
        "minimal": "pure white, blue accent — clean & corporate",
    }
    # Build one line per theme and join them for a readable output
    lines = [f"  {name}: {desc}" for name, desc in descriptions.items()]
    return "Available themes:\n" + "\n".join(lines)


#entry point: when this script is run directly, start the MCP server.
# claude desktop launches this via the command/args in claude_desktop_config.json.
# the server communicates over stdio (MCP protocol) and blocks until terminated.
if __name__ == "__main__":
    mcp.run()   #start the MCP server - blocks indefinitely waiting for tool calls